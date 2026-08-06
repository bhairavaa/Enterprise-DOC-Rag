import docx
import pymupdf
from pptx import Presentation

from app.ingestion.loaders.docx_loader import load_docx
from app.ingestion.loaders.pdf_loader import load_pdf
from app.ingestion.loaders.pptx_loader import load_pptx


def test_load_pdf_tracks_last_seen_heading_as_section_title(tmp_path):
    path = tmp_path / "doc.pdf"
    pdf = pymupdf.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "Big Heading", fontsize=20)
    page.insert_text((72, 100), "Body text under the heading.", fontsize=10)
    pdf.save(str(path))
    pdf.close()

    documents = load_pdf(path)

    assert len(documents) == 1
    assert documents[0].metadata["page_label"] == "1"
    assert documents[0].metadata["section_title"] == "Big Heading"
    assert "Body text under the heading." in documents[0].text


def test_load_pdf_skips_blank_pages(tmp_path):
    path = tmp_path / "doc.pdf"
    pdf = pymupdf.open()
    pdf.new_page()  # blank
    page2 = pdf.new_page()
    page2.insert_text((72, 72), "Some content", fontsize=10)
    pdf.save(str(path))
    pdf.close()

    documents = load_pdf(path)

    assert len(documents) == 1
    assert documents[0].metadata["page_label"] == "2"


def test_load_docx_splits_on_headings(tmp_path):
    path = tmp_path / "doc.docx"
    doc = docx.Document()
    doc.add_heading("Doc Title", level=0)
    doc.add_heading("Section A", level=1)
    doc.add_paragraph("Content for section A.")
    doc.add_heading("Section B", level=1)
    doc.add_paragraph("Content for section B.")
    doc.save(str(path))

    documents = load_docx(path)

    assert [d.metadata["section_title"] for d in documents] == ["Section A", "Section B"]
    assert "Content for section A." in documents[0].text
    assert "Content for section B." in documents[1].text


def test_load_docx_without_headings_returns_single_document(tmp_path):
    path = tmp_path / "doc.docx"
    doc = docx.Document()
    doc.add_paragraph("Just a plain paragraph, no headings anywhere.")
    doc.save(str(path))

    documents = load_docx(path)

    assert len(documents) == 1
    assert documents[0].metadata["section_title"] is None


def test_load_pptx_extracts_title_body_and_notes_without_duplication(tmp_path):
    path = tmp_path / "deck.pptx"
    presentation = Presentation()
    layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(layout)
    slide.shapes.title.text = "Slide Title"
    body = slide.placeholders[1].text_frame
    body.text = "First bullet"
    body.add_paragraph().text = "Second bullet"
    slide.notes_slide.notes_text_frame.text = "Speaker notes here."
    presentation.save(str(path))

    documents = load_pptx(path)

    assert len(documents) == 1
    doc = documents[0]
    assert doc.metadata["slide_number"] == "1"
    assert doc.metadata["section_title"] == "Slide Title"
    assert doc.metadata["notes_text"] == "Speaker notes here."
    # Title must appear exactly once in the body text, not duplicated.
    assert doc.text.count("Slide Title") == 1
    assert "First bullet" in doc.text
    assert "Second bullet" in doc.text
