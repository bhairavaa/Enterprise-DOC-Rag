from pathlib import Path

from llama_index.core import Document
from pptx import Presentation


def load_pptx(file_path: Path) -> list[Document]:
    """One Document per slide (a slide is a natural chunk boundary — the
    downstream splitter never merges text across slides since each is a
    separate Document). Speaker notes are captured separately so they can be
    filtered in/out of retrieval via metadata."""
    presentation = Presentation(str(file_path))

    documents: list[Document] = []
    for slide_num, slide in enumerate(presentation.slides, start=1):
        title = _slide_title(slide)
        body_text = _slide_body_text(slide, title)
        notes_text = _slide_notes(slide)

        if not body_text and not notes_text:
            continue

        documents.append(
            Document(
                text=body_text,
                metadata={
                    "slide_number": str(slide_num),
                    "section_title": title,
                    "notes_text": notes_text,
                },
            )
        )
    return documents


def _slide_title(slide) -> str | None:
    title_shape = slide.shapes.title
    if title_shape is not None and title_shape.has_text_frame:
        text = title_shape.text_frame.text.strip()
        return text or None
    return None


def _slide_body_text(slide, title: str | None) -> str:
    # slide.shapes.title returns a fresh wrapper object on each access, so
    # `is`/`==` against an object pulled from iterating slide.shapes doesn't
    # reliably match even for the same underlying XML element — compare by
    # shape_id instead.
    title_shape = slide.shapes.title
    title_id = title_shape.shape_id if title_shape is not None else None

    parts: list[str] = []
    for shape in slide.shapes:
        if shape.shape_id == title_id:
            continue
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if text:
            parts.append(text)
    body = "\n".join(parts).strip()
    if title:
        return f"{title}\n{body}".strip()
    return body


def _slide_notes(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    return (slide.notes_slide.notes_text_frame.text or "").strip()
