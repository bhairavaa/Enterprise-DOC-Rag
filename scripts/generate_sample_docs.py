"""Generates a small multi-tenant sample corpus under sample_docs/ for
exercising ingestion locally: two tenants, mixed PDF/DOCX/PPTX, each with
real heading/section structure so chunking and citation metadata have
something meaningful to key off of.

Usage: python scripts/generate_sample_docs.py
"""

from pathlib import Path

import docx
import pymupdf
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent.parent / "sample_docs"


def write_pdf(path: Path, title: str, sections: list[tuple[str, list[str]]]) -> None:
    doc = pymupdf.open()
    page = doc.new_page()
    y = 72
    page.insert_text((72, y), title, fontsize=20, fontname="helv")
    y += 40

    for heading, paragraphs in sections:
        if y > 720:
            page = doc.new_page()
            y = 72
        page.insert_text((72, y), heading, fontsize=15, fontname="helv")
        y += 26
        for para in paragraphs:
            for line in _wrap(para, 90):
                if y > 740:
                    page = doc.new_page()
                    y = 72
                page.insert_text((72, y), line, fontsize=10, fontname="helv")
                y += 16
            y += 8
        y += 16

    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(path))
    doc.close()


def write_docx(path: Path, title: str, sections: list[tuple[str, list[str]]]) -> None:
    doc = docx.Document()
    doc.add_heading(title, level=0)
    for heading, paragraphs in sections:
        doc.add_heading(heading, level=1)
        for para in paragraphs:
            doc.add_paragraph(para)
    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(path))


def write_pptx(path: Path, title: str, slides: list[tuple[str, list[str], str]]) -> None:
    presentation = Presentation()
    title_layout = presentation.slide_layouts[0]
    bullet_layout = presentation.slide_layouts[1]

    title_slide = presentation.slides.add_slide(title_layout)
    title_slide.shapes.title.text = title

    for heading, bullets, notes in slides:
        slide = presentation.slides.add_slide(bullet_layout)
        slide.shapes.title.text = heading
        body = slide.placeholders[1].text_frame
        body.text = bullets[0]
        for bullet in bullets[1:]:
            p = body.add_paragraph()
            p.text = bullet
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(path))


def _wrap(text: str, width: int) -> list[str]:
    words = text.split()
    lines, current = [], ""
    for word in words:
        if len(current) + len(word) + 1 > width:
            lines.append(current)
            current = word
        else:
            current = f"{current} {word}".strip()
    if current:
        lines.append(current)
    return lines


def main() -> None:
    # --- acme / engineering / policies ---
    write_pdf(
        ROOT / "acme" / "engineering" / "policies" / "remote_work_policy.pdf",
        "Remote Work Policy",
        [
            (
                "1. Eligibility",
                [
                    "All full-time engineering employees are eligible for remote work "
                    "after completing their 90-day onboarding period. Managers may "
                    "approve exceptions on a case-by-case basis for new hires with "
                    "prior remote experience.",
                ],
            ),
            (
                "2. Equipment and Expenses",
                [
                    "The company provides a laptop, monitor, and a one-time home "
                    "office stipend of $500. Internet reimbursement is capped at "
                    "$50 per month and must be submitted through the standard "
                    "expense system with a receipt attached.",
                ],
            ),
            (
                "3. Core Hours",
                [
                    "Remote engineers are expected to be reachable between 10am and "
                    "4pm in their local timezone to support cross-team "
                    "collaboration and incident response rotations.",
                ],
            ),
        ],
    )

    write_docx(
        ROOT / "acme" / "engineering" / "policies" / "code_review_guidelines.docx",
        "Code Review Guidelines",
        [
            (
                "Purpose",
                [
                    "This document describes the expectations for pull request "
                    "reviews across the engineering organization, including "
                    "turnaround time and required approvals.",
                ],
            ),
            (
                "Review SLA",
                [
                    "Reviewers should provide first feedback within one business "
                    "day. Pull requests tagged urgent should be reviewed within "
                    "four hours during business hours.",
                ],
            ),
            (
                "Required Approvals",
                [
                    "Changes to shared libraries require two approvals, one of "
                    "which must be from the library's designated code owner. "
                    "Application-level changes require a single approval.",
                ],
            ),
        ],
    )

    write_pptx(
        ROOT / "acme" / "engineering" / "handbook" / "onboarding_overview.pptx",
        "Engineering Onboarding Overview",
        [
            (
                "Week 1: Setup",
                [
                    "Provision laptop and accounts",
                    "Complete security training",
                    "Meet your onboarding buddy",
                ],
                "Make sure IT tickets are filed before the new hire's start date.",
            ),
            (
                "Week 2: First Contributions",
                [
                    "Ship a small bug fix",
                    "Attend architecture overview session",
                    "Shadow an on-call rotation",
                ],
                "",
            ),
        ],
    )

    # --- globex / hr / benefits ---
    write_pdf(
        ROOT / "globex" / "hr" / "benefits" / "health_benefits_guide.pdf",
        "Health Benefits Guide",
        [
            (
                "Plan Options",
                [
                    "Globex offers three health plan tiers: Basic, Standard, and "
                    "Premium. Employees can switch tiers once per year during "
                    "open enrollment in November.",
                ],
            ),
            (
                "Dependent Coverage",
                [
                    "Spouses and children up to age 26 can be added to any plan "
                    "tier. Proof of dependent status must be submitted within 30 "
                    "days of the qualifying life event.",
                ],
            ),
        ],
    )

    write_docx(
        ROOT / "globex" / "hr" / "benefits" / "pto_policy.docx",
        "Paid Time Off Policy",
        [
            (
                "Accrual",
                [
                    "Full-time employees accrue 1.5 days of PTO per month, up to "
                    "a maximum of 18 days per year, prorated for part-time staff.",
                ],
            ),
            (
                "Carryover",
                [
                    "Up to 5 unused PTO days may be carried over into the next "
                    "calendar year. Any balance beyond that is forfeited unless "
                    "local law requires otherwise.",
                ],
            ),
        ],
    )

    print(f"Sample corpus written under {ROOT}")


if __name__ == "__main__":
    main()
